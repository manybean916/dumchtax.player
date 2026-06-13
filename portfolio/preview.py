"""Convert dumchtax-portfolio.pptx → preview.html (single-file browser preview)"""
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import base64, io, json, re, html as html_mod

def rgb_from_color(color_format):
    try:
        if color_format and color_format.type:
            rgb = color_format.rgb
            return f"#{rgb}"
    except:
        pass
    return None

prs = Presentation("dumchtax-portfolio.pptx")
slides_data = []

for slide in prs.slides:
    elements = []
    sw = prs.slide_width.inches
    sh = prs.slide_height.inches

    for shape in slide.shapes:
        el = {"type": shape.shape_type, "x": shape.left.inches if shape.left else 0,
              "y": shape.top.inches if shape.top else 0,
              "w": shape.width.inches if shape.width else 0,
              "h": shape.height.inches if shape.height else 0}

        # Fill color for shapes
        if hasattr(shape, 'fill'):
            try:
                fg = shape.fill.fore_color
                if fg and fg.type:
                    el["fill"] = f"#{fg.rgb}"
            except:
                pass

        # Text
        if shape.has_text_frame:
            parts = []
            for para in shape.text_frame.paragraphs:
                para_parts = []
                for run in para.runs:
                    txt = run.text
                    if not txt:
                        continue
                    style = {}
                    try:
                        style["size"] = run.font.size.pt if run.font.size else None
                    except: pass
                    try:
                        style["bold"] = run.font.bold
                    except: pass
                    try:
                        style["italic"] = run.font.italic
                    except: pass
                    try:
                        c = run.font.color
                        if c and c.type:
                            style["color"] = f"#{c.rgb}"
                    except: pass
                    para_parts.append({"text": txt, "style": style})
                if para_parts:
                    parts.append(para_parts)
            if parts:
                el["text_parts"] = parts

        elements.append(el)

    # Get background color
    bg = "0A0A0A"
    try:
        bg_color = slide.background.fill.fore_color.rgb
        bg = str(bg_color)
    except:
        pass

    slides_data.append({"bg": bg, "elements": elements, "sw": sw, "sh": sh})

# Build HTML
html_slides = []
for i, slide in enumerate(slides_data):
    pct_w = lambda x: round(x / slide["sw"] * 100, 3)
    pct_h = lambda y: round(y / slide["sh"] * 100, 3)

    els_html = []
    for el in slide["elements"]:
        x = pct_w(el["x"])
        y = pct_h(el["y"])
        w = pct_w(el["w"])
        h = pct_h(el["h"])
        style_parts = [f"position:absolute;left:{x}%;top:{y}%;width:{w}%;height:{h}%;overflow:hidden"]

        fill = el.get("fill")
        if fill:
            style_parts.append(f"background:{fill}")

        if "text_parts" in el:
            inner_lines = []
            for para in el["text_parts"]:
                line_parts = []
                for run in para:
                    s = run["style"]
                    span_style = []
                    if s.get("size"):
                        # Scale font size relative to slide height (7.5" = 100%)
                        fs_pct = s["size"] / (slide["sh"] * 72) * 100
                        span_style.append(f"font-size:{fs_pct:.2f}cqh")
                    if s.get("bold"):
                        span_style.append("font-weight:700")
                    if s.get("italic"):
                        span_style.append("font-style:italic")
                    if s.get("color"):
                        span_style.append(f"color:{s['color']}")
                    txt = html_mod.escape(run["text"])
                    if span_style:
                        line_parts.append(f'<span style="{";".join(span_style)}">{txt}</span>')
                    else:
                        line_parts.append(txt)
                inner_lines.append("".join(line_parts))
            inner = "<br>".join(inner_lines)
            style_str = ";".join(style_parts)
            els_html.append(f'<div style="{style_str};font-family:Consolas,monospace;color:#F5F5F0;white-space:pre-wrap;line-height:1.2;container-type:size">{inner}</div>')
        elif fill:
            style_str = ";".join(style_parts)
            els_html.append(f'<div style="{style_str}"></div>')

    html_slides.append(f"""
<div class="slide" id="s{i+1}" style="background:#{slide['bg']}">
  <div class="slide-num">{i+1} / {len(slides_data)}</div>
  {''.join(els_html)}
</div>""")

html = f"""<!DOCTYPE html>
<html lang="ko">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width,initial-scale=1">
<title>Dumchtax Player — Portfolio Preview</title>
<style>
* {{ box-sizing:border-box; margin:0; padding:0 }}
body {{ background:#111; font-family:Consolas,monospace; color:#f5f5f0 }}
.toolbar {{ position:fixed; top:0; left:0; right:0; z-index:100;
  background:#0a0a0a; border-bottom:1px solid #222;
  display:flex; align-items:center; gap:12px; padding:8px 16px }}
.toolbar button {{ background:#913439; border:none; color:#fff; padding:6px 14px;
  font-family:Consolas; font-size:11px; cursor:pointer; letter-spacing:2px }}
.toolbar button:hover {{ background:#b04448 }}
.toolbar span {{ font-size:11px; color:#9a9a9a; letter-spacing:2px }}
#slide-counter {{ color:#f5f5f0; font-size:12px; letter-spacing:2px; margin-left:auto }}

.slides-container {{ padding-top:48px; display:flex; flex-direction:column; gap:12px; align-items:center; padding-bottom:40px }}
.slide {{ position:relative; width:90vw; max-width:1200px;
  aspect-ratio:16/9; overflow:hidden;
  box-shadow:0 8px 32px rgba(0,0,0,.6) }}
.slide-num {{ position:absolute; bottom:6px; right:10px;
  font-size:9px; color:#333; letter-spacing:2px; z-index:10; font-family:Consolas }}

/* Fullscreen mode */
body.fullscreen .slides-container {{ padding:0; gap:0 }}
body.fullscreen .slide {{ display:none; width:100vw; max-width:none;
  height:calc(100vh - 48px) }}
body.fullscreen .slide.active {{ display:block }}
body.fullscreen .slide .slide-num {{ font-size:10px }}
</style>
</head>
<body>
<div class="toolbar">
  <button onclick="prevSlide()">◀ PREV</button>
  <button onclick="nextSlide()">NEXT ▶</button>
  <button onclick="toggleView()">FULLSCREEN</button>
  <span id="slide-counter">1 / {len(slides_data)}</span>
  <span style="margin-left:auto;color:#5a5a5a">DUMCHTAX PLAYER — PORTFOLIO</span>
</div>
<div class="slides-container" id="container">
{''.join(html_slides)}
</div>
<script>
let current = 1;
const total = {len(slides_data)};
let fs = false;

function updateCounter() {{
  document.getElementById('slide-counter').textContent = current + ' / ' + total;
}}
function showSlide(n) {{
  current = Math.max(1, Math.min(total, n));
  if (fs) {{
    document.querySelectorAll('.slide').forEach((s,i) => {{
      s.classList.toggle('active', i+1 === current);
    }});
    document.getElementById('s'+current).scrollIntoView({{behavior:'instant'}});
  }} else {{
    document.getElementById('s'+current).scrollIntoView({{behavior:'smooth', block:'center'}});
  }}
  updateCounter();
}}
function prevSlide() {{ showSlide(current - 1); }}
function nextSlide() {{ showSlide(current + 1); }}
function toggleView() {{
  fs = !fs;
  document.body.classList.toggle('fullscreen', fs);
  if (fs) {{ showSlide(current); }}
}}
document.addEventListener('keydown', e => {{
  if (e.key==='ArrowRight'||e.key==='ArrowDown') nextSlide();
  if (e.key==='ArrowLeft'||e.key==='ArrowUp') prevSlide();
  if (e.key==='f'||e.key==='F') toggleView();
}});

// IntersectionObserver for scroll tracking
const obs = new IntersectionObserver(entries => {{
  entries.forEach(e => {{
    if (e.isIntersecting && e.intersectionRatio > 0.5) {{
      const m = e.target.id.match(/s(\\d+)/);
      if (m) {{ current = parseInt(m[1]); updateCounter(); }}
    }}
  }});
}}, {{threshold: 0.5}});
document.querySelectorAll('.slide').forEach(s => obs.observe(s));
</script>
</body>
</html>"""

with open("preview.html", "w") as f:
    f.write(html)
print("✓ preview.html saved")
